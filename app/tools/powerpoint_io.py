"""Build a proposal .pptx from a dict of {section: text}."""
from __future__ import annotations
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt


def build_proposal_pptx(
    sections: dict[str, str],
    title: str,
    client_name: str,
) -> bytes:
    """Title slide + one content slide per section. Returns .pptx bytes."""
    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = f"Prepared for {client_name}"

    # Content slides — one per section
    content_layout = prs.slide_layouts[1]
    for section_name, section_text in sections.items():
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = section_name
        tf = slide.placeholders[1].text_frame
        tf.text = section_text
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()