"""PDF, Office, plain text, and Markdown extractors.

Each extractor reads a source file and returns an ExtractedDocument with
per-unit text (pages for PDFs, sheets for spreadsheets, slides for decks,
sections for Word, one unit for plain text/Markdown) plus document-level
metadata. The Chunker (Step 2c) consumes ExtractedDocument without caring
about the original format.

Extensibility: a new file type adds a new class implementing the Extractor
protocol and registers it in EXTRACTORS_BY_EXT. No changes to the pipeline.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol, runtime_checkable

import openpyxl
import pdfplumber
from docx import Document as DocxDocument
from markdown_it import MarkdownIt
from pptx import Presentation

from app.observability.logging import get_logger

log = get_logger(__name__)


# ---------- Data shapes returned by extractors --------------------------------


@dataclass
class ExtractedUnit:
    """One logical unit of an extracted document.

    For PDFs, a unit is a page. For Excel, a unit is a worksheet.
    For PowerPoint, a unit is a slide. For Word/MD/text, the whole document
    is typically one unit (chunking is the chunker's job, not ours).
    The chunker treats units as the outermost split boundary.
    """

    index: int
    label: str
    text: str


@dataclass
class ExtractedDocument:
    """A whole document, broken into ordered units, with metadata."""

    source_path: str
    source_format: str  # "pdf", "xlsx", "docx", "pptx", "txt", "md", "code"
    page_count: int
    units: list[ExtractedUnit] = field(default_factory=list)
    metadata: dict[str, str] = field(default_factory=dict)

    @property
    def full_text(self) -> str:
        """All units joined with double newlines. Useful for whole-doc embeddings."""
        return "\n\n".join(u.text for u in self.units if u.text.strip())


# ---------- The Protocol every extractor implements ---------------------------


@runtime_checkable
class Extractor(Protocol):
    """Anything callable with .extract(path) returning ExtractedDocument."""

    def extract(self, path: Path) -> ExtractedDocument: ...


# ---------- Concrete extractors -----------------------------------------------


class PdfExtractor:
    """Extract text from a PDF, one ExtractedUnit per page."""

    def extract(self, path: Path) -> ExtractedDocument:
        units: list[ExtractedUnit] = []
        metadata: dict[str, str] = {}

        with pdfplumber.open(path) as pdf:
            metadata = {
                k: str(v) for k, v in (pdf.metadata or {}).items() if v is not None
            }

            for i, page in enumerate(pdf.pages):
                page_text = (page.extract_text() or "").strip()

                tables = page.extract_tables() or []
                if tables:
                    blocks = []
                    for t in tables:
                        rows = [
                            "\t".join((cell or "").strip() for cell in row)
                            for row in t
                            if any(cell for cell in row)
                        ]
                        if rows:
                            blocks.append("\n".join(rows))
                    if blocks:
                        page_text = (page_text + "\n\n" + "\n\n".join(blocks)).strip()

                if not page_text:
                    log.warning(
                        "extractor.empty_page",
                        path=str(path),
                        page_index=i,
                        hint="likely scanned/image-only PDF; OCR needed",
                    )

                units.append(
                    ExtractedUnit(index=i, label=f"page {i + 1}", text=page_text)
                )

            page_count = len(pdf.pages)

        log.info(
            "extractor.pdf_done",
            path=str(path),
            pages=page_count,
            empty_pages=sum(1 for u in units if not u.text),
        )

        return ExtractedDocument(
            source_path=str(path),
            source_format="pdf",
            page_count=page_count,
            units=units,
            metadata=metadata,
        )


class XlsxExtractor:
    """Extract data from an Excel workbook, one ExtractedUnit per worksheet."""

    def extract(self, path: Path) -> ExtractedDocument:
        wb = openpyxl.load_workbook(filename=path, read_only=True, data_only=True)
        units: list[ExtractedUnit] = []

        for i, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                if any(c is not None for c in row):
                    rows.append("\t".join("" if c is None else str(c) for c in row))
            units.append(
                ExtractedUnit(index=i, label=sheet_name, text="\n".join(rows))
            )

        wb.close()

        log.info("extractor.xlsx_done", path=str(path), sheets=len(units))

        return ExtractedDocument(
            source_path=str(path),
            source_format="xlsx",
            page_count=len(units),
            units=units,
        )


class DocxExtractor:
    """Extract text from a Word document.

    Returns the whole document as one unit. The chunker handles splitting.
    Captures headings, paragraphs, and tables. Headers/footers are skipped
    (usually boilerplate).
    """

    def extract(self, path: Path) -> ExtractedDocument:
        doc = DocxDocument(str(path))

        parts: list[str] = []

        # Body: paragraphs and tables in document order
        for element in doc.element.body.iter():
            tag = element.tag.split("}")[-1]
            if tag == "p":
                # Paragraph: collect text from all runs
                text = "".join(t.text or "" for t in element.iter() if t.tag.endswith("}t"))
                if text.strip():
                    parts.append(text)
            elif tag == "tbl":
                # Table: serialize as TSV
                rows = []
                for tr in element.iter():
                    if tr.tag.split("}")[-1] == "tr":
                        cells = []
                        for tc in tr.iter():
                            if tc.tag.split("}")[-1] == "tc":
                                cell_text = "".join(
                                    t.text or ""
                                    for t in tc.iter()
                                    if t.tag.endswith("}t")
                                )
                                cells.append(cell_text.strip())
                        if any(cells):
                            rows.append("\t".join(cells))
                if rows:
                    parts.append("\n".join(rows))

        text = "\n\n".join(parts)

        metadata: dict[str, str] = {}
        if doc.core_properties.title:
            metadata["title"] = doc.core_properties.title
        if doc.core_properties.author:
            metadata["author"] = doc.core_properties.author

        log.info("extractor.docx_done", path=str(path), char_count=len(text))

        return ExtractedDocument(
            source_path=str(path),
            source_format="docx",
            page_count=1,
            units=[ExtractedUnit(index=0, label="document", text=text)],
            metadata=metadata,
        )


class PptxExtractor:
    """Extract text from a PowerPoint deck, one ExtractedUnit per slide.

    Captures title, body text, and speaker notes. Tables are serialized as
    TSV inline with slide text. Images are ignored (would need OCR).
    """

    def extract(self, path: Path) -> ExtractedDocument:
        prs = Presentation(str(path))
        units: list[ExtractedUnit] = []

        for i, slide in enumerate(prs.slides):
            parts: list[str] = []

            for shape in slide.shapes:
                # Text frames (titles, bullets, text boxes)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            parts.append(line)
                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            parts.append("\t".join(cells))

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[notes] {notes}")

            units.append(
                ExtractedUnit(
                    index=i, label=f"slide {i + 1}", text="\n".join(parts)
                )
            )

        log.info("extractor.pptx_done", path=str(path), slides=len(units))

        return ExtractedDocument(
            source_path=str(path),
            source_format="pptx",
            page_count=len(units),
            units=units,
        )


class PlainTextExtractor:
    """Extract any plain-text file. Used for .txt and (with language metadata) code.

    For code files, the file extension is captured in metadata["language"]
    so Phase 3's RAG can filter on language without language-aware parsing.
    """

    def extract(self, path: Path) -> ExtractedDocument:
        text = path.read_text(encoding="utf-8", errors="replace")

        ext = path.suffix.lower()
        source_format = "code" if ext in CODE_EXTENSIONS else "txt"
        metadata: dict[str, str] = {}
        if source_format == "code":
            metadata["language"] = CODE_EXTENSIONS[ext]

        log.info(
            "extractor.text_done",
            path=str(path),
            char_count=len(text),
            format=source_format,
        )

        return ExtractedDocument(
            source_path=str(path),
            source_format=source_format,
            page_count=1,
            units=[ExtractedUnit(index=0, label="document", text=text)],
            metadata=metadata,
        )


class MarkdownExtractor:
    """Extract text from a Markdown file.

    Uses markdown-it-py to tokenize. We render to plain text by walking
    the token stream and emitting text content, preserving structure
    (headings stay on their own line, lists keep their bullets).
    """

    _md = MarkdownIt()

    def extract(self, path: Path) -> ExtractedDocument:
        raw = path.read_text(encoding="utf-8", errors="replace")
        tokens = self._md.parse(raw)

        parts: list[str] = []
        for tok in tokens:
            if tok.type == "inline" and tok.content:
                parts.append(tok.content)
            elif tok.type == "code_block" or tok.type == "fence":
                if tok.content.strip():
                    parts.append(tok.content.strip())

        text = "\n\n".join(parts)

        log.info(
            "extractor.markdown_done",
            path=str(path),
            char_count=len(text),
        )

        return ExtractedDocument(
            source_path=str(path),
            source_format="md",
            page_count=1,
            units=[ExtractedUnit(index=0, label="document", text=text)],
        )


# ---------- Extension dispatch ------------------------------------------------

# Code file extensions and their language tag for retrieval metadata.
# Treated as plain text; no language-aware parsing.
CODE_EXTENSIONS: dict[str, str] = {
    ".py": "python",
    ".js": "javascript",
    ".ts": "typescript",
    ".java": "java",
    ".go": "go",
    ".rs": "rust",
    ".c": "c",
    ".cpp": "cpp",
    ".h": "c",
    ".hpp": "cpp",
    ".cs": "csharp",
    ".rb": "ruby",
    ".php": "php",
    ".sh": "shell",
    ".sql": "sql",
    ".yaml": "yaml",
    ".yml": "yaml",
    ".json": "json",
    ".xml": "xml",
    ".html": "html",
    ".css": "css",
    ".scala": "scala",
    ".kt": "kotlin",
    ".swift": "swift",
    ".r": "r",
    ".tf": "terraform",
}

EXTRACTORS_BY_EXT: dict[str, Extractor] = {
    # Documents
    ".pdf": PdfExtractor(),
    ".docx": DocxExtractor(),
    # Spreadsheets
    ".xlsx": XlsxExtractor(),
    ".xls": XlsxExtractor(),
    # Presentations
    ".pptx": PptxExtractor(),
    # Plain text
    ".txt": PlainTextExtractor(),
    ".md": MarkdownExtractor(),
    ".markdown": MarkdownExtractor(),
}

# Code extensions all use PlainTextExtractor.
for code_ext in CODE_EXTENSIONS:
    EXTRACTORS_BY_EXT.setdefault(code_ext, PlainTextExtractor())


class UnsupportedFormatError(ValueError):
    """Raised when no extractor is registered for a given file extension."""


def get_extractor(path: Path) -> Extractor:
    """Return the right extractor for the file's extension."""
    ext = path.suffix.lower()
    if ext not in EXTRACTORS_BY_EXT:
        raise UnsupportedFormatError(
            f"No extractor for {ext!r}. Supported: {sorted(EXTRACTORS_BY_EXT)}"
        )
    return EXTRACTORS_BY_EXT[ext]


def supported_extensions() -> list[str]:
    """Sorted list of supported file extensions. Useful for UI / API surfaces."""
    return sorted(EXTRACTORS_BY_EXT)


# ---------- CLI demo ----------------------------------------------------------


def main() -> None:
    """Quick CLI: `python -m app.ingestion.extractor path/to/file.<ext>`."""
    import sys

    if len(sys.argv) != 2:
        print("usage: python -m app.ingestion.extractor <path>")
        print(f"supported: {supported_extensions()}")
        sys.exit(2)

    path = Path(sys.argv[1])
    if not path.exists():
        print(f"file not found: {path}")
        sys.exit(1)

    extractor = get_extractor(path)
    doc = extractor.extract(path)

    print(f"format:     {doc.source_format}")
    print(f"units:      {len(doc.units)}")
    print(f"total text: {sum(len(u.text) for u in doc.units)} chars")
    print(f"metadata:   {doc.metadata}")
    print()
    for u in doc.units[:3]:
        preview = u.text[:200].replace("\n", " ⏎ ")
        print(f"  [{u.label}] {preview}{'...' if len(u.text) > 200 else ''}")
    if len(doc.units) > 3:
        print(f"  ... and {len(doc.units) - 3} more units")


if __name__ == "__main__":
    main()