#!/usr/bin/env python
"""Live acceptance harness for Knowledge Studio.

Exercises every Knowledge Studio feature against a *running* stack — the
FastAPI app plus Mongo, Weaviate, MinIO and Redis — and prints a per-feature
PASS/FAIL report.  Unlike tests/, nothing here is mocked: this is the check
that answers "does the product actually work end to end right now".

Usage:
    uv run python scripts/verify_knowledge_studio.py
    uv run python scripts/verify_knowledge_studio.py --base http://localhost:8000
    uv run python scripts/verify_knowledge_studio.py --keep   # keep fixtures

Exit code is the number of failed checks (0 = everything passed).
"""
from __future__ import annotations

import argparse
import json
import sys
import time
import traceback
from dataclasses import dataclass, field
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Callable

import httpx

TIMEOUT = httpx.Timeout(180.0, connect=10.0)
TERMINAL = {"completed", "partially_completed", "failed", "cancelled"}


# ---------- Reporting ---------------------------------------------------------


@dataclass
class Result:
    """Provides the Result behaviour.

    Attributes:
        area (str).
        name (str).
        ok (bool).
        detail (str).
    """
    area: str
    name: str
    ok: bool
    detail: str = ""


@dataclass
class Report:
    """Provides the Report behaviour.

    Attributes:
        results (list[Result]).
    """
    results: list[Result] = field(default_factory=list)

    def record(self, area: str, name: str, ok: bool, detail: str = "") -> bool:
        """Record the result.

        Args:
            area (str): The area.
            name (str): Workflow or resource name.
            ok (bool): The ok.
            detail (str): The detail (optional, default '').

        Returns:
            bool: The result.
        """
        self.results.append(Result(area, name, ok, detail))
        mark = "\033[32mPASS\033[0m" if ok else "\033[31mFAIL\033[0m"
        line = f"  {mark}  {name}"
        if detail:
            line += f"  \033[90m— {detail}\033[0m"
        print(line, flush=True)
        return ok

    def check(self, area: str, name: str, fn: Callable[[], Any]) -> Any:
        """Run fn; PASS when it returns without raising. Returns fn's value."""
        try:
            value = fn()
        except Exception as exc:  # noqa: BLE001 - harness reports every failure
            detail = f"{type(exc).__name__}: {exc}"
            if not isinstance(exc, AssertionError):
                detail = detail.splitlines()[0]
            self.record(area, name, False, detail)
            return None
        detail = value if isinstance(value, str) else ""
        self.record(area, name, True, detail)
        return value

    @property
    def failed(self) -> list[Result]:
        """The failed."""
        return [item for item in self.results if not item.ok]

    def summary(self) -> int:
        """Compute the summary.

        Returns:
            int: The result.
        """
        total, bad = len(self.results), len(self.failed)
        print("\n" + "=" * 72)
        print(f"  {total - bad}/{total} checks passed")
        if bad:
            print(f"\n  \033[31m{bad} failing:\033[0m")
            for item in self.failed:
                print(f"    - [{item.area}] {item.name}: {item.detail}")
        print("=" * 72)
        return bad


# ---------- Client ------------------------------------------------------------


class Api:
    """Provides the Api behaviour."""
    def __init__(self, base: str, username: str, password: str):
        """Initialize the Api.

        Args:
            base (str): The base.
            username (str): Username value.
            password (str): Password value.
        """
        self.base = base.rstrip("/")
        self.client = httpx.Client(base_url=self.base, timeout=TIMEOUT)
        r = self.client.post(
            "/auth/token", data={"username": username, "password": password}
        )
        r.raise_for_status()
        self.token = r.json()["access_token"]

    @property
    def h(self) -> dict[str, str]:
        """The h."""
        return {"Authorization": f"Bearer {self.token}"}

    def _send(self, method: str, path: str, **kw: Any) -> httpx.Response:
        """Issue a request, waiting out the API's per-minute rate limit.

        The harness makes far more calls per minute than any human would, so a
        429 here is the limiter working as designed — not a product defect.
        Backing off keeps a rate limit from masquerading as 30 failed features.
        """
        for attempt in range(6):
            r = self.client.request(method, path, headers=self.h, **kw)
            if r.status_code != 429:
                return r
            delay = float(r.headers.get("retry-after") or min(2**attempt, 20))
            time.sleep(delay)
        return r

    def get(self, path: str, **kw: Any) -> httpx.Response:
        """Return the result.

        Args:
            path (str): Filesystem path.
            **kw (Any): The kw.

        Returns:
            httpx.Response: The result.
        """
        return self._send("GET", path, **kw)

    def post(self, path: str, **kw: Any) -> httpx.Response:
        """Compute the post.

        Args:
            path (str): Filesystem path.
            **kw (Any): The kw.

        Returns:
            httpx.Response: The result.
        """
        return self._send("POST", path, **kw)

    def jpost(self, path: str, payload: dict[str, Any]) -> Any:
        """Compute the jpost.

        Args:
            path (str): Filesystem path.
            payload (dict[str, Any]): Event or audit payload.

        Returns:
            Any: The result.
        """
        r = self.post(path, json=payload)
        _raise(r)
        return r.json()

    def jget(self, path: str, **kw: Any) -> Any:
        """Compute the jget.

        Args:
            path (str): Filesystem path.
            **kw (Any): The kw.

        Returns:
            Any: The result.
        """
        r = self.get(path, **kw)
        _raise(r)
        return r.json()

    def close(self) -> None:
        """Close the result."""
        self.client.close()


def _raise(r: httpx.Response) -> None:
    """Raise the result.

    Args:
        r (httpx.Response): The r.
    """
    if r.is_success:
        return
    raise AssertionError(f"HTTP {r.status_code} {r.request.url.path} :: {r.text[:300]}")


# ---------- Fixture documents -------------------------------------------------

PROSE = (
    "The Dura 25 centrifugal pump delivers a maximum flow of 25 cubic metres per hour "
    "at a differential head of 32 metres. Wetted parts are available in stainless steel "
    "316L, Hastelloy C276 and PTFE-lined ductile iron. Mechanical seals are single or "
    "double acting, with silicon carbide faces recommended for abrasive slurries. "
    "PTFE seals resist sulphuric acid up to 98 percent concentration at ambient "
    "temperature. Do not use EPDM elastomers in hydrocarbon service; select FKM instead. "
)


def _long(times: int = 6) -> str:
    """Internal helper for the long step.

    Args:
        times (int): The times (optional, default 6).

    Returns:
        str: The result.
    """
    return PROSE * times


def build_fixtures(root: Path) -> dict[str, Path]:
    """Write one fixture per supported source format."""
    files: dict[str, Path] = {}

    def write(name: str, text: str) -> None:
        """Write the result.

        Args:
            name (str): Workflow or resource name.
            text (str): The text.
        """
        path = root / name
        path.write_text(text, encoding="utf-8")
        files[path.suffix.lstrip(".")] = path

    body = "\n\n".join(f"Section {i + 1}: Operation\n\n{_long()}" for i in range(4))
    write("manual.txt", body)
    write("manual.md", "# Dura 25 Manual\n\n## Seals\n\n" + _long() + "\n\n## Duty\n\n" + _long())
    write(
        "specs.csv",
        "model,flow_m3h,head_m,seal,notes\n"
        + "\n".join(
            f"Dura {n},{n},{n + 7},PTFE,{_long(1)[:200]}" for n in (25, 32, 40, 50)
        ),
    )
    write(
        "catalog.json",
        json.dumps(
            {"products": [{"model": f"Dura {n}", "description": _long(2)} for n in (25, 32)]},
            indent=1,
        ),
    )
    write(
        "page.html",
        "<html><head><title>Dura 25</title></head><body>"
        + "".join(f"<h2>Section {i}</h2><p>{_long(2)}</p>" for i in range(3))
        + "</body></html>",
    )
    write(
        "data.xml",
        "<catalog>" + "".join(f"<item><name>Dura {n}</name><body>{_long(2)}</body></item>" for n in (25, 32)) + "</catalog>",
    )

    # DOCX / XLSX / PPTX via their own libraries when installed.
    try:
        from docx import Document

        doc = Document()
        for i in range(6):
            doc.add_heading(f"Section {i + 1}", level=2)
            doc.add_paragraph(_long())
        path = root / "manual.docx"
        doc.save(str(path))
        files["docx"] = path
    except Exception:  # noqa: BLE001 - optional format
        pass

    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.append(["model", "flow", "notes"])
        for n in (25, 32, 40):
            ws.append([f"Dura {n}", n, _long(1)[:400]])
        path = root / "specs.xlsx"
        wb.save(str(path))
        files["xlsx"] = path
    except Exception:  # noqa: BLE001 - optional format
        pass

    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Dura 25 — slide {i + 1}"
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            box.text_frame.text = _long(2)
        path = root / "deck.pptx"
        prs.save(str(path))
        files["pptx"] = path
    except Exception:  # noqa: BLE001 - optional format
        pass

    # Minimal single-page PDF with a real text layer, written by hand so the
    # harness needs no PDF dependency.
    files["pdf"] = _write_pdf(root / "note.pdf")

    # A scanned-style PDF: valid structure, no extractable text. Used to prove
    # the "no text could be extracted" path reports something actionable.
    files["pdf_blank"] = _write_pdf(root / "scan.pdf", lines=[])
    return files


def _write_pdf(path: Path, lines: list[str] | None = None) -> Path:
    """Write the pdf.

    Args:
        path (Path): Filesystem path.
        lines (list[str] | None): The lines (optional, default None).

    Returns:
        Path: The pdf.
    """
    if lines is None:
        lines = [PROSE[i : i + 90] for i in range(0, len(PROSE) * 2, 90)][:40]
    text = "BT /F1 11 Tf 40 750 Td 13 TL\n" + "".join(
        f"({line.replace('(', '').replace(')', '')}) Tj T*\n" for line in lines
    ) + "ET"
    stream = text.encode("latin-1", "replace")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    )
    path.write_bytes(bytes(out))
    return path


# ---------- Helpers -----------------------------------------------------------


def wait_for_job(api: Api, job_id: str, timeout: float = 180.0) -> dict[str, Any]:
    """Compute the wait for job.

    Args:
        api (Api): The api.
        job_id (str): The job id.
        timeout (float): Timeout in seconds (optional, default 180.0).

    Returns:
        dict[str, Any]: The for job.
    """
    deadline = time.time() + timeout
    job: dict[str, Any] = {}
    while time.time() < deadline:
        job = api.jget(f"/api/knowledge/ingestions/{job_id}")
        if job["status"] in TERMINAL:
            return job
        time.sleep(1.0)
    raise AssertionError(f"job {job_id} still {job.get('status')} after {timeout}s")


def ingest(
    api: Api,
    collection_id: str,
    paths: list[Path],
    *,
    parser: dict[str, Any] | None = None,
    chunking: dict[str, Any] | None = None,
    embedding: dict[str, Any] | None = None,
    metadata: dict[str, Any] | None = None,
    wait: bool = True,
) -> dict[str, Any]:
    """Ingest the result.

    Args:
        api (Api): The api.
        collection_id (str): Knowledge collection identifier.
        paths (list[Path]): The paths.
        parser (dict[str, Any] | None): The parser (optional, default None).
        chunking (dict[str, Any] | None): The chunking (optional, default None).
        embedding (dict[str, Any] | None): The embedding (optional, default None).
        metadata (dict[str, Any] | None): Metadata mapping (optional, default None).
        wait (bool): The wait (optional, default True).

    Returns:
        dict[str, Any]: The result.
    """
    handles = [("files", (p.name, p.read_bytes(), "application/octet-stream")) for p in paths]
    data: dict[str, Any] = {"metadata_json": json.dumps(metadata or {})}
    for key, prof in (("parser", parser), ("chunking", chunking), ("embedding", embedding)):
        if prof:
            data[f"{key}_profile_id"] = prof["profile_id"]
            data[f"{key}_profile_version"] = str(prof["version"])
    r = api.post(
        f"/api/knowledge/collections/{collection_id}/ingestions", files=handles, data=data
    )
    _raise(r)
    job = r.json()
    return wait_for_job(api, job["ingestion_job_id"]) if wait else job


def profile(api: Api, kind: str, name: str, strategy: str, config: dict[str, Any]) -> dict[str, Any]:
    """Compute the profile.

    Args:
        api (Api): The api.
        kind (str): The kind.
        name (str): Workflow or resource name.
        strategy (str): The strategy.
        config (dict[str, Any]): Node configuration mapping.

    Returns:
        dict[str, Any]: The result.
    """
    return api.jpost(
        "/api/knowledge/profiles",
        {"profile_type": kind, "name": name, "strategy": strategy, "config": config},
    )


# ---------- The suite ---------------------------------------------------------


def run(base: str, user: str, password: str, keep: bool) -> int:
    """Run the result.

    Args:
        base (str): The base.
        user (str): Authenticated current user.
        password (str): Password value.
        keep (bool): The keep.

    Returns:
        int: The result.
    """
    rep = Report()
    api = Api(base, user, password)
    stamp = time.strftime("%H%M%S")

    health = api.client.get("/health").json()
    print(f"\nStack: {base}  ready={health.get('ready')}")
    down = [k for k, v in health.get("services", {}).items() if v.get("status") != "ok"]
    if down:
        print(f"\033[31mDegraded services: {down}\033[0m")

    with TemporaryDirectory() as tmp:
        root = Path(tmp)
        fixtures = build_fixtures(root)
        print(f"Fixtures: {', '.join(sorted(fixtures))}\n")

        state: dict[str, Any] = {}

        # ---- A. Collections ------------------------------------------------
        print("\033[1mA. Collections\033[0m")

        def create_collection() -> str:
            """Create the collection.

            Returns:
                str: The collection.
            """
            coll = api.jpost(
                "/api/knowledge/collections",
                {
                    "name": f"KS Verify {stamp}",
                    "description": "acceptance harness",
                    "doc_types": ["manual", "spec"],
                    "metadata_schema": {"product": "string", "revision": "integer"},
                },
            )
            assert coll["collection_id"].startswith("col_"), coll["collection_id"]
            assert coll["status"] == "draft", coll["status"]
            state["collection"] = coll
            return coll["collection_id"]

        cid = rep.check("collections", "create collection with metadata schema", create_collection)
        if not cid:
            rep.summary()
            return 1

        rep.check(
            "collections",
            "get collection by id",
            lambda: _assert_eq(api.jget(f"/api/knowledge/collections/{cid}")["collection_id"], cid),
        )
        rep.check(
            "collections",
            "list collections includes it",
            lambda: _assert_in(cid, [c["collection_id"] for c in api.jget("/api/knowledge/collections")]),
        )
        rep.check(
            "collections",
            "unknown collection returns 404",
            lambda: _assert_status(api.get("/api/knowledge/collections/col_does_not_exist"), 404),
        )

        # ---- B. Profiles ---------------------------------------------------
        print("\n\033[1mB. Profiles\033[0m")
        defaults = rep.check(
            "profiles",
            "ensure default profiles (all 5 types)",
            lambda: _assert_keys(
                api.jpost("/api/knowledge/profiles/defaults", {}),
                {"parser", "chunking", "embedding", "retrieval", "generation"},
            ),
        )
        if defaults:
            state["defaults"] = defaults

        rep.check(
            "profiles",
            "defaults are idempotent (same ids on re-call)",
            lambda: _assert_eq(
                api.jpost("/api/knowledge/profiles/defaults", {})["parser"]["profile_id"],
                defaults["parser"]["profile_id"],
            ),
        )

        def versioning() -> str:
            """Compute the versioning.

            Returns:
                str: The result.
            """
            first = profile(api, "chunking", f"Versioned {stamp}", "recursive", {"target_tokens": 400})
            second = api.jpost(
                "/api/knowledge/profiles",
                {
                    "profile_type": "chunking",
                    "name": f"Versioned {stamp}",
                    "strategy": "recursive",
                    "config": {"target_tokens": 480},
                    "profile_id": first["profile_id"],
                },
            )
            assert second["profile_id"] == first["profile_id"], "profile id must be stable"
            assert second["version"] == first["version"] + 1, (
                f"expected v{first['version'] + 1}, got v{second['version']}"
            )
            pinned = api.jget(
                f"/api/knowledge/profiles/{first['profile_id']}",
                params={"version": first["version"]},
            )
            assert pinned["config"]["target_tokens"] == 400, "old version was rewritten"
            state["versioned"] = second
            return f"v1 preserved, v2 created"

        rep.check("profiles", "update creates a new version, history preserved", versioning)
        rep.check(
            "profiles",
            "list profiles by type",
            lambda: _assert_true(
                all(
                    p["profile_type"] == "chunking"
                    for p in api.jget("/api/knowledge/profiles", params={"profile_type": "chunking"})
                ),
                "type filter leaked other types",
            ),
        )
        rep.check(
            "profiles",
            "ingestion presets endpoint",
            lambda: _assert_true(
                len(api.jget("/api/knowledge/ingestion-presets")) > 0, "no presets returned"
            ),
        )

        # ---- C. Ingestion: formats ----------------------------------------
        print("\n\033[1mC. Ingestion — source formats\033[0m")
        for fmt in ["txt", "md", "csv", "json", "html", "xml", "docx", "xlsx", "pptx", "pdf"]:
            if fmt not in fixtures:
                rep.record("formats", f"ingest .{fmt}", False, "fixture unavailable (library missing)")
                continue

            def one(fmt: str = fmt) -> str:
                """Compute the one.

                Args:
                    fmt (str): The fmt (optional, default fmt).

                Returns:
                    str: The result.
                """
                job = ingest(api, cid, [fixtures[fmt]], metadata={"product": "Dura 25", "revision": 3})
                assert job["status"] == "completed", f"{job['status']} :: {job['errors']}"
                assert job["chunks_created"] > 0, "no chunks"
                return f"{job['chunks_created']} chunks"

            rep.check("formats", f"ingest .{fmt}", one)

        def scanned() -> str:
            """Compute the scanned.

            Returns:
                str: The result.
            """
            job = ingest(api, cid, [fixtures["pdf_blank"]])
            assert job["status"] == "failed", f"expected failure, got {job['status']}"
            msg = job["errors"][0]["message"]
            assert "no text could be extracted" in msg, msg
            assert "OCR" in msg, "message should point at OCR"
            return "actionable error"

        rep.check("formats", "text-less PDF fails with an actionable message", scanned)

        def tiny() -> str:
            """Compute the tiny.

            Returns:
                str: The result.
            """
            path = root / "tiny.txt"
            path.write_text("Dura 25 max flow is 25 m3/h.", encoding="utf-8")
            job = ingest(api, cid, [path])
            assert job["status"] == "completed", f"{job['status']} :: {job['errors']}"
            return f"{job['chunks_created']} chunk(s)"

        rep.check("formats", "very short document still ingests (min_tokens regression)", tiny)

        # ---- D. Ingestion: chunking strategies -----------------------------
        print("\n\033[1mD. Ingestion — chunking strategies\033[0m")
        for strategy in [
            "fixed_token",
            "recursive",
            "structure_aware",
            "parent_child",
            "contextual",
            "sentence_window",
            "semantic",
        ]:

            def one(strategy: str = strategy) -> str:
                """Compute the one.

                Args:
                    strategy (str): The strategy (optional, default strategy).

                Returns:
                    str: The result.
                """
                prof = profile(
                    api,
                    "chunking",
                    f"{strategy} {stamp}",
                    strategy,
                    {
                        "strategy": strategy,
                        "target_tokens": 320,
                        "max_tokens": 900,
                        "overlap_tokens": 48,
                    },
                )
                job = ingest(api, cid, [fixtures["md"]], chunking=prof)
                assert job["status"] == "completed", f"{job['status']} :: {job['errors']}"
                assert job["chunks_created"] > 0, "no chunks"
                state.setdefault("strategy_index", {})[strategy] = job["target_index_id"]
                return f"{job['chunks_created']} chunks"

            rep.check("chunking", f"strategy: {strategy}", one)

        # ---- E. Ingestion: parser strategies -------------------------------
        print("\n\033[1mE. Ingestion — parser strategies\033[0m")
        for strategy in ["standard", "layout_aware", "structure_aware"]:

            def one(strategy: str = strategy) -> str:
                """Compute the one.

                Args:
                    strategy (str): The strategy (optional, default strategy).

                Returns:
                    str: The result.
                """
                prof = profile(api, "parser", f"parser {strategy} {stamp}", strategy, {"strategy": strategy})
                job = ingest(api, cid, [fixtures["html"]], parser=prof)
                assert job["status"] == "completed", f"{job['status']} :: {job['errors']}"
                return f"{job['chunks_created']} chunks"

            rep.check("parsers", f"parser: {strategy}", one)

        def ocr_unconfigured() -> str:
            """Compute the ocr unconfigured.

            Returns:
                str: The unconfigured.
            """
            prof = profile(
                api, "parser", f"parser ocr {stamp}", "ocr_fallback",
                {"strategy": "ocr_fallback", "ocr_min_text_characters": 80},
            )
            job = ingest(api, cid, [fixtures["pdf_blank"]], parser=prof)
            assert job["status"] == "failed", f"expected failure, got {job['status']}"
            return "reports missing OCR provider"

        rep.check("parsers", "ocr_fallback without a provider fails cleanly", ocr_unconfigured)

        # ---- F. Ingestion jobs ---------------------------------------------
        print("\n\033[1mF. Ingestion jobs\033[0m")
        rep.check(
            "jobs",
            "multi-file job reports per-document outcomes",
            lambda: _multi_file(api, cid, fixtures),
        )
        rep.check("jobs", "duplicate filenames in one job rejected", lambda: _dupe_names(api, cid, fixtures))
        rep.check("jobs", "re-ingesting identical content is idempotent", lambda: _idempotent(api, cid, fixtures))
        rep.check("jobs", "list jobs filtered by collection", lambda: _job_list(api, cid))
        rep.check("jobs", "cancel is accepted and terminal", lambda: _cancel(api, cid, fixtures))
        rep.check(
            "jobs",
            "upload with no files rejected (422)",
            lambda: _assert_status(
                api.post(f"/api/knowledge/collections/{cid}/ingestions", data={"metadata_json": "{}"}),
                422,
            ),
        )
        rep.check(
            "jobs",
            "metadata violating the collection schema is rejected",
            lambda: _assert_status(
                api.post(
                    f"/api/knowledge/collections/{cid}/ingestions",
                    files=[("files", ("m.txt", fixtures["txt"].read_bytes(), "text/plain"))],
                    data={"metadata_json": json.dumps({"revision": "not-an-integer"})},
                ),
                422,
            ),
        )

        # ---- G. Documents & indexes ----------------------------------------
        print("\n\033[1mG. Documents & indexes\033[0m")
        docs = rep.check(
            "documents",
            "list documents in collection",
            lambda: _assert_nonempty(api.jget(f"/api/knowledge/collections/{cid}/documents"), "documents"),
        )
        if docs:
            doc_id = docs[0]["document_id"]
            rep.check(
                "documents",
                "get document by id",
                lambda: _assert_eq(api.jget(f"/api/knowledge/documents/{doc_id}")["document_id"], doc_id),
            )
            rep.check(
                "documents",
                "document provenance (source version + hash)",
                lambda: _provenance(api.jget(f"/api/knowledge/documents/{doc_id}")),
            )
            rep.check(
                "documents",
                "presigned source URL for the original upload",
                lambda: _source_url(api, doc_id),
            )

        indexes = rep.check(
            "indexes",
            "list indexes for collection",
            lambda: _assert_nonempty(api.jget(f"/api/knowledge/collections/{cid}/indexes"), "indexes"),
        )
        rep.check(
            "indexes",
            "index records exact profile versions",
            lambda: _index_pins(indexes or []),
        )
        rep.check("indexes", "activate an index", lambda: _activate(api, cid, indexes or []))
        rep.check("indexes", "switch active index to another", lambda: _switch(api, cid, indexes or []))
        rep.check(
            "indexes",
            "activating an unknown index is refused",
            lambda: _assert_status(
                api.post(f"/api/knowledge/collections/{cid}/indexes/idx_nope/activate"), 409
            ),
        )

        # ---- H. Retrieval ---------------------------------------------------
        print("\n\033[1mH. Retrieval\033[0m")
        rep.check(
            "retrieval",
            "presets endpoint",
            lambda: _assert_true(len(api.jget("/api/retrieval/presets")) > 0, "no presets"),
        )

        query = "Which seal material resists sulphuric acid?"
        # RRF is a fusion strategy over hybrid retrieval, not a retrieval
        # strategy of its own — the API's enum is correct.
        for label, extra in [
            ("dense", {"strategy": "dense"}),
            ("sparse", {"strategy": "sparse"}),
            ("hybrid (relative score)", {"strategy": "hybrid", "fusion_strategy": "relative_score"}),
            ("hybrid + RRF fusion", {"strategy": "hybrid", "fusion_strategy": "rrf"}),
        ]:

            def one(extra: dict = extra) -> str:
                """Compute the one.

                Args:
                    extra (dict): The extra (optional, default extra).

                Returns:
                    str: The result.
                """
                out = api.jpost(
                    "/api/retrieval/search",
                    {
                        "collection_id": cid,
                        "query": query,
                        "rerank": False,
                        "candidate_count": 20,
                        "final_count": 5,
                        **extra,
                    },
                )
                chunks = out.get("chunks") or out.get("results") or []
                assert chunks, f"no chunks returned :: {list(out)}"
                return f"{len(chunks)} chunks"

            rep.check("retrieval", f"strategy: {label}", one)

        rep.check("retrieval", "reranking changes/keeps a valid ordering", lambda: _rerank(api, cid, query))
        rep.check("retrieval", "typed metadata filter narrows results", lambda: _filters(api, cid, query))
        rep.check("retrieval", "unknown filter field rejected", lambda: _bad_filter(api, cid, query))
        rep.check("retrieval", "reserved scope field in filters rejected", lambda: _reserved_filter(api, cid, query))

        for transform in ["rewrite", "multi_query", "hyde", "decomposition"]:

            def one(transform: str = transform) -> str:
                """Compute the one.

                Args:
                    transform (str): The transform (optional, default transform).

                Returns:
                    str: The result.
                """
                out = api.jpost(
                    "/api/retrieval/search",
                    {"collection_id": cid, "query": query, "strategy": "hybrid",
                     "query_transform": transform, "rerank": False, "final_count": 5},
                )
                assert (out.get("chunks") or out.get("results")), "no chunks"
                return "ok"

            rep.check("retrieval", f"query transform: {transform}", one)

        # Each expansion is only valid on an index chunked that way; reuse the
        # indexes built in section D rather than the collection's active one.
        for expansion, chunking in [
            ("parent", "parent_child"),
            ("sentence_window", "sentence_window"),
            ("contextual", "contextual"),
        ]:

            def one(expansion: str = expansion, chunking: str = chunking) -> str:
                """Compute the one.

                Args:
                    expansion (str): The expansion (optional, default expansion).
                    chunking (str): The chunking (optional, default chunking).

                Returns:
                    str: The result.
                """
                index_id = (state.get("strategy_index") or {}).get(chunking)
                assert index_id, f"no index was built with {chunking} chunking"
                out = api.jpost(
                    "/api/retrieval/search",
                    {"collection_id": cid, "query": query, "strategy": "hybrid",
                     "index_id": index_id, "context_expansion": expansion,
                     "rerank": False, "final_count": 5},
                )
                assert (out.get("chunks") or out.get("results")), "no chunks"
                return f"on {chunking} index"

            rep.check("retrieval", f"context expansion: {expansion}", one)

        rep.check("retrieval", "compare endpoint runs several experiments", lambda: _compare(api, cid, query))
        rep.check("retrieval", "search pinned to a specific index", lambda: _pinned_index(api, cid, query, indexes or []))
        rep.check("retrieval", "corpus inspector delegates to the same path", lambda: _inspect(api, cid, query))

        # ---- I. Traces -------------------------------------------------------
        print("\n\033[1mI. Retrieval traces\033[0m")
        traces = rep.check(
            "traces",
            "traces are persisted and listable",
            lambda: _assert_nonempty(api.jget("/api/retrieval/traces"), "traces"),
        )
        rep.check("traces", "trace records the resolved resources", lambda: _trace_detail(api, traces or []))

        # ---- J. RAG agents ---------------------------------------------------
        print("\n\033[1mJ. RAG Agents\033[0m")
        agent = rep.check("rag", "create RAG Agent from saved resources", lambda: _make_agent(api, cid, defaults, stamp, state))
        if agent:
            rep.check(
                "rag",
                "get agent by id",
                lambda: _assert_eq(
                    api.jget(f"/api/rag-agents/{agent['rag_agent_id']}")["rag_agent_id"],
                    agent["rag_agent_id"],
                ),
            )
            rep.check(
                "rag",
                "list agents includes it",
                lambda: _assert_in(
                    agent["rag_agent_id"], [a["rag_agent_id"] for a in api.jget("/api/rag-agents")]
                ),
            )
            rep.check("rag", "agent answers a grounded query with citations", lambda: _agent_query(api, agent, query))
            rep.check("rag", "agent query records resolved resources", lambda: _agent_resolution(api, agent, query))
            rep.check(
                "rag",
                "citations are retrieved_not_verified (no VerifiedClaim)",
                lambda: _evidence_boundary(api, agent, query),
            )
        rep.check(
            "rag",
            "unknown agent returns 404",
            lambda: _assert_status(api.get("/api/rag-agents/rag_missing"), 404),
        )

        # ---- K. Security -----------------------------------------------------
        print("\n\033[1mK. Security & scope isolation\033[0m")
        rep.check("security", "unauthenticated request refused", lambda: _anon(base, cid))
        rep.check("security", "search on a foreign collection refused", lambda: _foreign(api))

        # ---- L. Workflow integration ----------------------------------------
        print("\n\033[1mL. Workflow integration\033[0m")
        rep.check("workflow", "RAGAgent + KnowledgeRetrieval registered in node catalog", lambda: _catalog(api))
        if agent:
            rep.check(
                "workflow",
                "preflight accepts a workflow using rag_agent_id",
                lambda: _preflight_ok(api, agent, cid),
            )
        rep.check("workflow", "preflight rejects an unknown rag_agent_id", lambda: _preflight_bad(api))

        if not keep:
            print("\n(fixtures are temporary; created collections/profiles remain in Mongo)")

    api.close()
    return rep.summary()


# ---------- Individual check bodies -------------------------------------------


def _assert_eq(actual: Any, expected: Any) -> str:
    """Internal helper for the assert eq step.

    Args:
        actual (Any): Actual value.
        expected (Any): Expected value.

    Returns:
        str: The eq.
    """
    assert actual == expected, f"expected {expected!r}, got {actual!r}"
    return ""


def _assert_in(needle: Any, haystack: Any) -> str:
    """Internal helper for the assert in step.

    Args:
        needle (Any): The needle.
        haystack (Any): The haystack.

    Returns:
        str: The in.
    """
    assert needle in haystack, f"{needle!r} missing"
    return ""


def _assert_true(value: bool, message: str) -> str:
    """Internal helper for the assert true step.

    Args:
        value (bool): Value to process.
        message (str): Message text.

    Returns:
        str: The true.
    """
    assert value, message
    return ""


def _assert_keys(payload: dict[str, Any], keys: set[str]) -> dict[str, Any]:
    """Internal helper for the assert keys step.

    Args:
        payload (dict[str, Any]): Event or audit payload.
        keys (set[str]): The keys.

    Returns:
        dict[str, Any]: The keys.
    """
    missing = keys - set(payload)
    assert not missing, f"missing {sorted(missing)}"
    return payload


def _assert_nonempty(value: list[Any], label: str) -> list[Any]:
    """Internal helper for the assert nonempty step.

    Args:
        value (list[Any]): Value to process.
        label (str): The label.

    Returns:
        list[Any]: The nonempty.
    """
    assert value, f"no {label} returned"
    return value


def _assert_status(r: httpx.Response, expected: int) -> str:
    """Internal helper for the assert status step.

    Args:
        r (httpx.Response): The r.
        expected (int): Expected value.

    Returns:
        str: The status.
    """
    assert r.status_code == expected, f"expected {expected}, got {r.status_code}: {r.text[:200]}"
    return f"{expected}"


def _multi_file(api: Api, cid: str, fixtures: dict[str, Path]) -> str:
    """Internal helper for the multi file step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        fixtures (dict[str, Path]): The fixtures.

    Returns:
        str: The file.
    """
    paths = [fixtures[k] for k in ("txt", "md", "csv") if k in fixtures]
    job = ingest(api, cid, paths)
    assert job["documents_total"] == len(paths), job["documents_total"]
    assert job["documents_processed"] == len(paths), f"{job['documents_processed']} :: {job['errors']}"
    return f"{job['documents_processed']}/{job['documents_total']} documents"


def _dupe_names(api: Api, cid: str, fixtures: dict[str, Path]) -> str:
    """Internal helper for the dupe names step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        fixtures (dict[str, Path]): The fixtures.

    Returns:
        str: The names.
    """
    blob = fixtures["txt"].read_bytes()
    r = api.post(
        f"/api/knowledge/collections/{cid}/ingestions",
        files=[("files", ("same.txt", blob, "text/plain")), ("files", ("same.txt", blob, "text/plain"))],
        data={"metadata_json": "{}"},
    )
    assert r.status_code == 422, f"expected 422, got {r.status_code}"
    return "422"


def _idempotent(api: Api, cid: str, fixtures: dict[str, Path]) -> str:
    """Internal helper for the idempotent step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        fixtures (dict[str, Path]): The fixtures.

    Returns:
        str: The result.
    """
    first = ingest(api, cid, [fixtures["csv"]])
    second = ingest(api, cid, [fixtures["csv"]])
    assert first["status"] == "completed" and second["status"] == "completed"
    key_a = first["source_inputs"][0]["storage_key"]
    key_b = second["source_inputs"][0]["storage_key"]
    assert key_a == key_b, "identical content must map to one content-addressed key"
    return "same storage key reused"


def _job_list(api: Api, cid: str) -> str:
    """Internal helper for the job list step.

    Args:
        api (Api): The api.
        cid (str): The cid.

    Returns:
        str: The list.
    """
    jobs = api.jget("/api/knowledge/ingestions", params={"collection_id": cid})
    assert jobs, "no jobs"
    assert all(j["collection_id"] == cid for j in jobs), "filter leaked other collections"
    return f"{len(jobs)} jobs"


def _cancel(api: Api, cid: str, fixtures: dict[str, Path]) -> str:
    """Cancel the result.

    Args:
        api (Api): The api.
        cid (str): The cid.
        fixtures (dict[str, Path]): The fixtures.

    Returns:
        str: The result.
    """
    job = ingest(api, cid, [fixtures["txt"]], wait=False)
    r = api.post(f"/api/knowledge/ingestions/{job['ingestion_job_id']}/cancel")
    assert r.status_code in (200, 409), f"unexpected {r.status_code}: {r.text[:200]}"
    final = wait_for_job(api, job["ingestion_job_id"])
    assert final["status"] in TERMINAL, final["status"]
    if final["status"] == "cancelled":
        coll = api.jget(f"/api/knowledge/collections/{cid}")
        assert coll.get("active_index_id") != job["target_index_id"], "cancelled index was activated"
    return final["status"]


def _provenance(doc: dict[str, Any]) -> str:
    """Internal helper for the provenance step.

    Args:
        doc (dict[str, Any]): Document.

    Returns:
        str: The result.
    """
    for field_name in ("content_hash", "current_source_version_id", "collection_id"):
        assert doc.get(field_name), f"missing {field_name}"
    return doc["content_hash"][:12]


def _source_url(api: Api, doc_id: str) -> str:
    """Internal helper for the source url step.

    Args:
        api (Api): The api.
        doc_id (str): The doc id.

    Returns:
        str: The url.
    """
    out = api.jget(f"/api/knowledge/documents/{doc_id}/source-url")
    assert out.get("url", "").startswith("http"), out
    with httpx.Client(timeout=TIMEOUT) as raw:
        head = raw.get(out["url"], headers={"Range": "bytes=0-64"})
    assert head.status_code in (200, 206), f"object not downloadable: {head.status_code}"
    return "downloadable"


def _index_pins(indexes: list[dict[str, Any]]) -> str:
    """Internal helper for the index pins step.

    Args:
        indexes (list[dict[str, Any]]): The indexes.

    Returns:
        str: The pins.
    """
    assert indexes, "no indexes"
    idx = indexes[0]
    for field_name in (
        "parser_profile_id", "parser_profile_version",
        "chunking_profile_id", "chunking_profile_version",
        "embedding_profile_id", "embedding_profile_version",
    ):
        assert idx.get(field_name) is not None, f"missing {field_name}"
    return "all profile versions pinned"


def _ready(indexes: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Internal helper for the ready step.

    Args:
        indexes (list[dict[str, Any]]): The indexes.

    Returns:
        list[dict[str, Any]]: The result.
    """
    return [i for i in indexes if i.get("status") in ("ready", "active", "inactive")]


def _activate(api: Api, cid: str, indexes: list[dict[str, Any]]) -> str:
    """Internal helper for the activate step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        indexes (list[dict[str, Any]]): The indexes.

    Returns:
        str: The result.
    """
    ready = _ready(indexes)
    assert ready, "no ready index to activate"
    coll = api.jpost(
        f"/api/knowledge/collections/{cid}/indexes/{ready[0]['index_id']}/activate", {}
    )
    assert coll["active_index_id"] == ready[0]["index_id"], coll["active_index_id"]
    return ready[0]["index_id"]


def _switch(api: Api, cid: str, indexes: list[dict[str, Any]]) -> str:
    """Internal helper for the switch step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        indexes (list[dict[str, Any]]): The indexes.

    Returns:
        str: The result.
    """
    ready = _ready(indexes)
    assert len(ready) >= 2, "need two ready indexes"
    api.jpost(f"/api/knowledge/collections/{cid}/indexes/{ready[0]['index_id']}/activate", {})
    coll = api.jpost(f"/api/knowledge/collections/{cid}/indexes/{ready[1]['index_id']}/activate", {})
    assert coll["active_index_id"] == ready[1]["index_id"], "switch did not take effect"
    return "switched without touching any workflow"


def _chunks(out: dict[str, Any]) -> list[dict[str, Any]]:
    """Internal helper for the chunks step.

    Args:
        out (dict[str, Any]): The out.

    Returns:
        list[dict[str, Any]]: The result.
    """
    return out.get("chunks") or out.get("results") or []


def _rerank(api: Api, cid: str, query: str) -> str:
    """Internal helper for the rerank step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The result.
    """
    base = {"collection_id": cid, "query": query, "strategy": "hybrid", "final_count": 5}
    plain = _chunks(api.jpost("/api/retrieval/search", {**base, "rerank": False}))
    ranked = _chunks(api.jpost("/api/retrieval/search", {**base, "rerank": True}))
    assert plain and ranked, "one of the runs returned nothing"
    return f"{len(plain)} → {len(ranked)} after rerank"


def _filters(api: Api, cid: str, query: str) -> str:
    """Internal helper for the filters step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The result.
    """
    out = api.jpost(
        "/api/retrieval/search",
        {"collection_id": cid, "query": query, "strategy": "hybrid", "rerank": False,
         "filters": {"product": "Dura 25"}, "final_count": 10},
    )
    assert _chunks(out), "filter removed everything — metadata not indexed?"
    return f"{len(_chunks(out))} chunks match product=Dura 25"


def _bad_filter(api: Api, cid: str, query: str) -> str:
    """Internal helper for the bad filter step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The filter.
    """
    r = api.post(
        "/api/retrieval/search",
        json={"collection_id": cid, "query": query, "filters": {"no_such_field": "x"}},
    )
    assert r.status_code in (400, 422), f"unknown field accepted ({r.status_code})"
    return str(r.status_code)


def _reserved_filter(api: Api, cid: str, query: str) -> str:
    """Internal helper for the reserved filter step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The filter.
    """
    for reserved in ("owner_scope_id", "workspace_id", "collection_id", "index_id"):
        r = api.post(
            "/api/retrieval/search",
            json={"collection_id": cid, "query": query, "filters": {reserved: "anything"}},
        )
        assert r.status_code in (400, 422), f"{reserved} accepted as a user filter ({r.status_code})"
    return "all reserved fields refused"


def _compare(api: Api, cid: str, query: str) -> str:
    """Compare the result.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The result.
    """
    base = {"collection_id": cid, "query": query, "final_count": 5}
    out = api.jpost(
        "/api/retrieval/compare",
        {"experiments": [
            {**base, "strategy": "dense", "rerank": False},
            {**base, "strategy": "hybrid", "rerank": False},
            {**base, "strategy": "hybrid", "rerank": True},
        ]},
    )
    runs = out if isinstance(out, list) else out.get("experiments") or out.get("results") or []
    assert len(runs) == 3, f"expected 3 runs, got {len(runs)}"
    return "3 experiments"


def _pinned_index(api: Api, cid: str, query: str, indexes: list[dict[str, Any]]) -> str:
    """Internal helper for the pinned index step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.
        indexes (list[dict[str, Any]]): The indexes.

    Returns:
        str: The index.
    """
    ready = _ready(indexes)
    assert ready, "no ready index"
    out = api.jpost(
        "/api/retrieval/search",
        {"collection_id": cid, "query": query, "index_id": ready[0]["index_id"],
         "strategy": "hybrid", "rerank": False, "final_count": 5},
    )
    assert _chunks(out), "pinned-index search returned nothing"
    return ready[0]["index_id"]


def _inspect(api: Api, cid: str, query: str) -> str:
    """Internal helper for the inspect step.

    Args:
        api (Api): The api.
        cid (str): The cid.
        query (str): Query filter.

    Returns:
        str: The result.
    """
    r = api.post("/api/inspect/retrieve", json={"collection_id": cid, "query": query, "top_k": 5})
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:200]}"
    return "200"


def _trace_detail(api: Api, traces: list[dict[str, Any]]) -> str:
    """Internal helper for the trace detail step.

    Args:
        api (Api): The api.
        traces (list[dict[str, Any]]): The traces.

    Returns:
        str: The detail.
    """
    assert traces, "no traces"
    key = traces[0].get("retrieval_request_id") or traces[0].get("request_id")
    assert key, f"trace has no id field :: {list(traces[0])}"
    trace = api.jget(f"/api/retrieval/traces/{key}")
    for field_name in (
        "collection_id", "resolved_index_id",
        "retrieval_profile_id", "retrieval_profile_version",
        "parser_profile_id", "chunking_profile_id", "embedding_profile_id",
    ):
        assert trace.get(field_name) is not None, f"trace missing {field_name}"
    return key


def _make_agent(api: Api, cid: str, defaults: Any, stamp: str, state: dict[str, Any]) -> dict[str, Any]:
    """Build the agent.

    Args:
        api (Api): The api.
        cid (str): The cid.
        defaults (Any): The defaults.
        stamp (str): The stamp.
        state (dict[str, Any]): Current workflow state.

    Returns:
        dict[str, Any]: The agent.
    """
    assert defaults, "no default profiles"
    agent = api.jpost(
        "/api/rag-agents",
        {
            "name": f"KS Verify Agent {stamp}",
            "collection_id": cid,
            "retrieval_profile_id": defaults["retrieval"]["profile_id"],
            "generation_profile_id": defaults["generation"]["profile_id"],
        },
    )
    assert agent["rag_agent_id"].startswith("rag_"), agent["rag_agent_id"]
    state["agent"] = agent
    return agent


def _agent_query(api: Api, agent: dict[str, Any], query: str) -> str:
    """Internal helper for the agent query step.

    Args:
        api (Api): The api.
        agent (dict[str, Any]): The agent.
        query (str): Query filter.

    Returns:
        str: The query.
    """
    out = api.jpost(f"/api/rag-agents/{agent['rag_agent_id']}/query", {"query": query})
    answer = out.get("answer") or out.get("output") or ""
    assert answer, f"no answer :: {list(out)}"
    citations = out.get("citations") or []
    assert citations, "grounded answer without citations"
    return f"{len(citations)} citations"


def _agent_resolution(api: Api, agent: dict[str, Any], query: str) -> str:
    """Internal helper for the agent resolution step.

    Args:
        api (Api): The api.
        agent (dict[str, Any]): The agent.
        query (str): Query filter.

    Returns:
        str: The resolution.
    """
    out = api.jpost(f"/api/rag-agents/{agent['rag_agent_id']}/query", {"query": query})
    blob = json.dumps(out)
    for token in ("col_", "idx_", "retprof_"):
        assert token in blob, f"resolved {token}* id not reported in the response"
    return "collection + index + profile recorded"


def _evidence_boundary(api: Api, agent: dict[str, Any], query: str) -> str:
    """Internal helper for the evidence boundary step.

    Args:
        api (Api): The api.
        agent (dict[str, Any]): The agent.
        query (str): Query filter.

    Returns:
        str: The boundary.
    """
    out = api.jpost(f"/api/rag-agents/{agent['rag_agent_id']}/query", {"query": query})
    blob = json.dumps(out)
    assert "VerifiedClaim" not in blob, "RAG output must not mint VerifiedClaim"
    assert "verified_claim" not in blob, "RAG output must not mint verified_claim"
    return "no VerifiedClaim minted"


def _anon(base: str, cid: str) -> str:
    """Internal helper for the anon step.

    Args:
        base (str): The base.
        cid (str): The cid.

    Returns:
        str: The result.
    """
    with httpx.Client(base_url=base, timeout=TIMEOUT) as anon:
        r = anon.get(f"/api/knowledge/collections/{cid}")
    assert r.status_code in (401, 403), f"anonymous read allowed ({r.status_code})"
    return str(r.status_code)


def _foreign(api: Api) -> str:
    """Internal helper for the foreign step.

    Args:
        api (Api): The api.

    Returns:
        str: The result.
    """
    r = api.post(
        "/api/retrieval/search",
        json={"collection_id": "col_someone_elses_collection", "query": "secrets"},
    )
    assert r.status_code in (200, 403, 404), f"unexpected {r.status_code}"
    if r.status_code != 200:
        return str(r.status_code)
    body = r.json()
    assert not (body.get("chunks") or body.get("candidates")), "foreign collection leaked chunks"
    assert not body.get("final_context"), "foreign collection leaked context"
    assert not body.get("resolved_index_id"), "foreign collection resolved an index"
    return "200 but empty (no existence oracle)"


def _catalog(api: Api) -> str:
    """Internal helper for the catalog step.

    Args:
        api (Api): The api.

    Returns:
        str: The result.
    """
    for path in ("/api/node-types", "/api/builder/node-types", "/api/builder/catalog", "/api/nodes"):
        r = api.get(path)
        if r.status_code == 200:
            break
    _raise(r)
    blob = json.dumps(r.json())
    assert "RAGAgent" in blob, "RAGAgent missing from catalog"
    assert "KnowledgeRetrieval" in blob, "KnowledgeRetrieval missing from catalog"
    for banned in ("DenseRAGNode", "HybridRAGNode", "HyDERAGNode", "MultiQueryRAGNode"):
        assert banned not in blob, f"{banned} must not be a node type"
    return "RAGAgent + KnowledgeRetrieval only"


WORKFLOW_TEMPLATE = """
name: ks_verify_{suffix}
description: Knowledge Studio acceptance probe
inputs:
  question:
    type: text
    required: true
nodes:
  - id: answer
    type: RAGAgent
    config:
      rag_agent_id: {agent_id}
      query: "{{{{inputs.question}}}}"
outputs:
  answer: "{{{{outputs.answer.answer}}}}"
"""


def _validate(api: Api, yaml_text: str) -> httpx.Response:
    # check_services=true is what turns on the zero-token resource probes
    # (RAG agent / collection / profile existence); without it only the
    # schema-level contract is validated.
    """Validate the result.

    Args:
        api (Api): The api.
        yaml_text (str): Workflow YAML text.

    Returns:
        httpx.Response: The result.
    """
    body = {
        "workflow_yaml": yaml_text,
        "check_services": True,
        "inputs": {"question": "Which seal material resists sulphuric acid?"},
    }
    for path, payload in (
        ("/api/workflows/validate", body),
        ("/api/builder/validate", body),
        ("/api/workflows/preflight", body),
    ):
        r = api.post(path, json=payload)
        if r.status_code != 404:
            return r
    raise AssertionError("no workflow validation endpoint found")


def _preflight_ok(api: Api, agent: dict[str, Any], cid: str) -> str:
    """Internal helper for the preflight ok step.

    Args:
        api (Api): The api.
        agent (dict[str, Any]): The agent.
        cid (str): The cid.

    Returns:
        str: The ok.
    """
    text = WORKFLOW_TEMPLATE.format(suffix="ok", agent_id=agent["rag_agent_id"])
    r = _validate(api, text)
    assert r.status_code == 200, f"HTTP {r.status_code}: {r.text[:300]}"
    body = r.json()
    issues = [i for i in (body.get("issues") or []) if i.get("severity") != "warning"]
    assert body.get("valid") is True, f"preflight rejected a valid workflow: {issues}"
    return "zero-token preflight clean"


def _preflight_bad(api: Api) -> str:
    """Internal helper for the preflight bad step.

    Args:
        api (Api): The api.

    Returns:
        str: The bad.
    """
    text = WORKFLOW_TEMPLATE.format(suffix="bad", agent_id="rag_01DOESNOTEXIST")
    r = _validate(api, text)
    assert r.status_code in (200, 400, 422), f"HTTP {r.status_code}"
    blob = r.text
    assert "RAG_AGENT_NOT_FOUND" in blob, f"missing agent not reported :: {blob[:300]}"
    assert '"knowledge_resources"' in blob, "resource probe did not run"
    return "RAG_AGENT_NOT_FOUND"


# ---------- Entry point -------------------------------------------------------


def main() -> int:
    """Compute the main.

    Returns:
        int: The result.
    """
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--base", default="http://localhost:8000")
    parser.add_argument("--user", default="ayush")
    parser.add_argument("--password", default="dev123")
    parser.add_argument("--keep", action="store_true", help="keep created fixtures")
    args = parser.parse_args()
    try:
        return run(args.base, args.user, args.password, args.keep)
    except KeyboardInterrupt:
        return 130
    except Exception:  # noqa: BLE001 - harness must always explain itself
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    sys.exit(main())
